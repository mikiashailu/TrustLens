"""
Build TrustLens 20-slide pitch deck from brand guidelines + unified / AI-model doc themes.
Outputs: docs/TrustLens-Pitch-Deck.pptx

Run: pip install -r scripts/requirements-presentations.txt
     python scripts/build_trustlens_pitch_deck.py

Animations:
  - Each slide gets a slide transition (fade or push) via OOXML — visible in Slide Show.
  - For bullet-by-bullet or fly-in: open in PowerPoint, select the text box, Animations tab,
    choose an entrance effect, Effect Options → By paragraph (on click).
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# --- Brand (from app style) ---
NAVY = RGBColor(0x00, 0x0F, 0x22)
TEAL = RGBColor(0x44, 0xDD, 0xC2)
EMERALD = RGBColor(0x00, 0x6B, 0x5C)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
BORDER_GRAY = RGBColor(0xC4, 0xC6, 0xCE)
BODY_TEXT = RGBColor(0x1A, 0x23, 0x32)
SUBTLE = RGBColor(0x5C, 0x6B, 0x7F)

FONT_MAIN = "Inter"
FONT_FALLBACK = "Calibri"  # used if Inter missing on build machine; PPT will substitute

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _font(run) -> None:
    run.font.name = FONT_MAIN
    run.font._element.set(qn("a:latin"), FONT_MAIN)
    try:
        run.font._element.set(qn("a:cs"), FONT_MAIN)
    except Exception:
        pass


def add_slide_transition(slide, kind: str = "fade") -> None:
    """Append p:transition under p:sld (fade or push)."""
    el = slide._element
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    # Remove existing transition if re-running tooling on template
    for child in list(el):
        if child.tag == qn("p:transition"):
            el.remove(child)
    if kind == "push":
        xml = f'<p:transition xmlns:p="{ns}" spd="med"><p:push dir="l"/></p:transition>'
    else:
        xml = f'<p:transition xmlns:p="{ns}" spd="med"><p:fade/></p:transition>'
    el.append(parse_xml(xml))


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, NAVY)
    add_slide_transition(slide, "fade")

    # Teal accent bar
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(2.85), Inches(1.85), Pt(6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(11.5), Inches(1.35))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TrustLens AI"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _font(p.runs[0])

    sub = tf.add_paragraph()
    sub.text = "Identity · Trust scoring · Eligibility"
    sub.font.size = Pt(22)
    sub.font.color.rgb = TEAL
    _font(sub.runs[0])

    foot = slide.shapes.add_textbox(Inches(0.55), Inches(6.35), Inches(11), Inches(0.6))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Unified platform + AI model overview"
    fp.font.size = Pt(14)
    fp.font.color.rgb = RGBColor(0xB8, 0xC5, 0xD6)
    _font(fp.runs[0])


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    transition: str = "fade",
    callout: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    add_slide_transition(slide, transition)

    # Heading
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.85))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    _font(tp.runs[0])

    # Bullets
    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.45), Inches(11.8), Inches(5.2))
    bf = body.text_frame
    bf.word_wrap = True
    bf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = line
        para.level = 0
        para.font.size = Pt(18)
        para.font.color.rgb = BODY_TEXT
        para.space_after = Pt(10)
        _font(para.runs[0])

    if callout:
        cx, cy, cw, ch = Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.95)
        call = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
        call.fill.solid()
        call.fill.fore_color.rgb = EMERALD
        call.line.color.rgb = EMERALD
        tf = call.text_frame
        tf.margin_left = tf.margin_right = Pt(16)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = tf.paragraphs[0]
        cp.text = callout
        cp.font.size = Pt(15)
        cp.font.bold = True
        cp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cp.alignment = PP_ALIGN.CENTER
        _font(cp.runs[0])


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    add_slide_transition(slide, "push")

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.75))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Architecture (logical)"
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    _font(tp.runs[0])

    boxes = [
        ("Mobile & Admin", Inches(0.9), Inches(1.55), Inches(3.4), Inches(1.05), False),
        ("FastAPI API", Inches(4.85), Inches(1.55), Inches(3.4), Inches(1.05), True),
        ("Trust pipeline\nOCR · rules", Inches(8.8), Inches(1.55), Inches(3.4), Inches(1.05), False),
        ("PostgreSQL", Inches(2.5), Inches(3.35), Inches(3.2), Inches(0.95), False),
        ("File storage", Inches(6.1), Inches(3.35), Inches(3.2), Inches(0.95), False),
    ]
    shapes_list = []
    for label, left, top, w, h, teal_node in boxes:
        sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        sh.line.color.rgb = TEAL if teal_node else BORDER_GRAY
        sh.line.width = Pt(2.5) if teal_node else Pt(1.5)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(15)
        p.font.color.rgb = NAVY
        p.font.bold = True
        _font(p.runs[0])
        shapes_list.append(sh)

    # Emerald arrows (connect approx centers)
    def arrow(x1, y1, x2, y2):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        c.line.color.rgb = EMERALD
        c.line.width = Pt(2.25)

    arrow(Inches(4.3), Inches(2.08), Inches(4.85), Inches(2.08))
    arrow(Inches(8.25), Inches(2.08), Inches(8.8), Inches(2.08))
    arrow(Inches(10.0), Inches(2.6), Inches(7.2), Inches(3.35))
    arrow(Inches(5.7), Inches(2.6), Inches(4.5), Inches(3.35))

    note = slide.shapes.add_textbox(Inches(0.55), Inches(4.85), Inches(12), Inches(1.8))
    nf = note.text_frame
    np = nf.paragraphs[0]
    np.text = (
        "Main flow: clients (JWT) → API → services → PostgreSQL + uploads. "
        "Teal nodes = user-facing entry; emerald arrows = primary request path."
    )
    np.font.size = Pt(14)
    np.font.color.rgb = SUBTLE
    _font(np.runs[0])


def add_scoring_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    add_slide_transition(slide, "fade")

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.75))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Scoring pipeline"
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    _font(tp.runs[0])

    # Positions in inches (float) for reliable connector math
    layout = [
        ("Inputs", 0.55, 1.4, 1.85, 0.75),
        ("Features", 2.55, 1.4, 1.85, 0.75),
        ("Criteria\n0–1", 4.55, 1.4, 1.85, 0.75),
        ("Modalities\n0–100", 6.55, 1.4, 1.95, 0.75),
        ("Combined", 8.75, 1.4, 1.85, 0.75),
        ("Eligibility", 10.75, 1.4, 1.85, 0.75),
    ]
    y_mid = Inches(1.4 + 0.75 / 2)
    for i, (text, x, y, w, h) in enumerate(layout):
        left, top, wi, he = Inches(x), Inches(y), Inches(w), Inches(h)
        teal_highlight = i in (0, len(layout) - 1)
        sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, wi, he)
        sh.fill.solid()
        sh.fill.fore_color.rgb = TEAL if teal_highlight else RGBColor(0xFF, 0xFF, 0xFF)
        sh.line.color.rgb = TEAL if teal_highlight else BORDER_GRAY
        sh.line.width = Pt(1.5)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
        _font(p.runs[0])
        if i > 0:
            _, x0, _, w0, _ = layout[i - 1]
            x1 = x
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x0 + w0),
                y_mid,
                Inches(x1),
                y_mid,
            )
            conn.line.color.rgb = EMERALD
            conn.line.width = Pt(2)


def add_eligibility_table_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    add_slide_transition(slide, "fade")

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.75))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Eligibility bands (combined score)"
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    _font(tp.runs[0])

    rows = [
        ["Score", "Loan", "Device ≥30", "Card ≥50"],
        ["≤ 25", "No", "No", "No"],
        ["26 – 50", "1–5k", "Yes", "No"],
        ["51 – 69", "5k–10k", "Yes", "Yes"],
        ["≥ 70", "10k–150k", "Yes", "Yes"],
    ]
    left0, top0, col_w, row_h = Inches(0.75), Inches(1.45), Inches(2.85), Inches(0.62)
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            left = left0 + c * col_w
            top = top0 + r * row_h
            sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, col_w - Inches(0.06), row_h - Inches(0.05))
            sh.fill.solid()
            sh.fill.fore_color.rgb = NAVY if r == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            sh.line.color.rgb = BORDER_GRAY
            tf = sh.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = cell
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14 if r == 0 else 13)
            p.font.bold = r == 0
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else BODY_TEXT
            _font(p.runs[0])


def add_thank_you(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)
    add_slide_transition(slide, "fade")
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(3.2), Inches(1.85), Pt(6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.9), Inches(11.5), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _font(p.runs[0])
    sub = tb.text_frame.add_paragraph()
    sub.text = "Questions?"
    sub.font.size = Pt(24)
    sub.font.color.rgb = TEAL
    _font(sub.runs[0])


def main() -> int:
    root = Path(__file__).resolve().parents[1]
    out = root / "docs" / "TrustLens-Pitch-Deck.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- 20 slides ---
    add_title_slide(prs)

    add_content_slide(
        prs,
        "Agenda",
        [
            "Problem & vision (slides 3–4)",
            "Platform overview & architecture (5–8)",
            "AI model: modalities & scoring (9–14)",
            "Tech stack & Trust Card (15–17)",
            "Strengths, roadmap, Q&A (18–20)",
        ],
        transition="push",
        callout="≈10 min walkthrough · demo-ready",
    )

    add_content_slide(
        prs,
        "The challenge",
        [
            "Digital finance needs identity + trust signals without opaque black boxes.",
            "eKYC produces documents, video, and audio — hard to fuse into one defensible score.",
            "Product teams need eligibility messaging that auditors and users can understand.",
        ],
        callout="Explainability is the product feature",
    )

    add_content_slide(
        prs,
        "TrustLens vision",
        [
            "Multi-modal trust: document + video + audio in one pipeline.",
            "Every check: pass / fail / uncertain with human-readable detail.",
            "Same API powers mobile users and admin operations.",
        ],
    )

    add_content_slide(
        prs,
        "What we built (system)",
        [
            "Mobile app: capture ID, video, audio → upload → scores & eligibility.",
            "Admin portal: dashboards, review, Trust Card workflows (same API).",
            "Backend: FastAPI, PostgreSQL, filesystem uploads, JWT-oriented security story.",
        ],
    )

    add_architecture_slide(prs)

    add_content_slide(
        prs,
        "Key API flows",
        [
            "Health: GET /health, GET /status (ops).",
            "Identity: POST /identity (multipart) → stored paths + submission row.",
            "Trust: POST /trust-result → full criterion breakdown + combined score.",
            "Products: POST /eligible → loan / device / card messaging from rules.",
        ],
    )

    add_content_slide(
        prs,
        "Trust Card (demo UX)",
        [
            "Digital credential when combined score > 45 (trust_card_service).",
            "POST /trust-card/issue, GET /trust-card, POST /trust-card/select product.",
            "Not a payment card — demo artifact tied to authenticated user.",
        ],
        callout="JWT-protected in target architecture",
    )

    add_scoring_flow_slide(prs)

    add_content_slide(
        prs,
        "What “AI” means here",
        [
            "No single end-to-end neural net for the headline score.",
            "Tesseract OCR = mature ML for ID text; fusion = averages + rules.",
            "Face / liveness / ASR slots exist as uncertain + guidance until models land.",
        ],
    )

    add_content_slide(
        prs,
        "Document modality",
        [
            "Pillow: resolution & clarity heuristics per ID side.",
            "Tesseract + PyMuPDF: OCR; match name, phone, sex, DOB, nationality to profile.",
            "Ethiopia-oriented phone patterns where applicable; English OCR default.",
        ],
    )

    add_content_slide(
        prs,
        "Video & audio modalities",
        [
            "Video: OpenCV for resolution & duration; file size as richness proxy.",
            "Audio: mutagen for duration & bitrate; size fallback.",
            "Future: face match, liveness, voice / ASR — same response shape.",
        ],
    )

    add_content_slide(
        prs,
        "Combined score (formulas)",
        [
            "Per modality: average of criterion scores in [0,1] × 100 → clamp 0–100.",
            "Combined: mean(document, video, audio) → clamp 0–100.",
            "Authoritative output: POST /trust-result (not upload placeholder).",
        ],
    )

    add_eligibility_table_slide(prs)

    add_content_slide(
        prs,
        "Modality balance",
        [
            "Eligibility response includes min, max, spread, weakest / strongest modality.",
            "Explains imbalance: e.g. great video, weak document.",
            "Helps UX and risk talk track remediation (re-upload, lighting, length).",
        ],
    )

    add_content_slide(
        prs,
        "Technologies (scoring stack)",
        [
            "FastAPI · Uvicorn · Pydantic · SQLAlchemy · PostgreSQL",
            "Tesseract · PyMuPDF · Pillow · OpenCV · mutagen · difflib",
            "Docker Compose for reproducible OCR + DB + volumes",
        ],
    )

    add_content_slide(
        prs,
        "Why this architecture wins",
        [
            "One scoring implementation — mobile and admin stay consistent.",
            "Criteria-level audit trail for demos, judges, and future compliance work.",
            "Clear extension points for real ML without breaking clients.",
        ],
        callout="Ship the demo · design for production",
    )

    add_content_slide(
        prs,
        "Roadmap (next)",
        [
            "Face embedding + active liveness on video.",
            "ASR for spoken-name consistency.",
            "Stronger fraud / document authenticity signals.",
            "Calibration with labeled operational data.",
        ],
    )

    add_content_slide(
        prs,
        "Demo checklist",
        [
            "Sign-in → profile → upload ID + AV → trust-result → eligible.",
            "Show per-check pass / fail / uncertain in UI.",
            "Optional: Trust Card issue + stats dashboard for admins.",
        ],
        transition="push",
    )

    add_thank_you(prs)

    prs.save(str(out))
    print(f"Wrote {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
